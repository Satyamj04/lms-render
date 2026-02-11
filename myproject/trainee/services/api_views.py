"""
Complete REST API for LMS Trainee Module
Handles courses, modules, media, quizzes, and assessments
"""
from rest_framework import viewsets, status
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import AllowAny, IsAuthenticated
from rest_framework.response import Response
from django.shortcuts import get_object_or_404
from django.db.models import Q, Count, Prefetch
from trainee.models import (
    Course, Module, VideoUnit, AudioUnit, PresentationUnit, 
    Quiz, Question, Assignment, Enrollment,
    UserProgress, ModuleCompletion, MediaMetadata, ScormPackage,
    TextUnit, PageUnit, Survey
)
from admin.models import UserProfile
import logging

logger = logging.getLogger(__name__)


# Helper function to get or create user session
def get_session_user(request):
    """Get UserProfile from request or return first available user"""
    from admin.models import UserProfile
    
    # Debug: Log what we're receiving
    logger.info(f"🔍 get_session_user called")
    logger.info(f"   request.data: {request.data if hasattr(request, 'data') else 'No data'}")
    logger.info(f"   request.GET: {dict(request.GET)}")
    
    # Check for user_id in request data or GET params FIRST (from frontend localStorage)
    user_id = request.data.get('user_id') if hasattr(request, 'data') else None
    if not user_id:
        user_id = request.GET.get('user_id')
    if not user_id:
        user_id = request.session.get('user_id')
    
    logger.info(f"   Extracted user_id: {user_id}")
    
    if user_id:
        try:
            user = UserProfile.objects.get(id=user_id)
            logger.info(f"✅ Found user by ID: {user.email} ({user_id})")
            return user
        except UserProfile.DoesNotExist:
            logger.warning(f"⚠️ User {user_id} not found in database")
    
    # Try Django authenticated user
    if request.user.is_authenticated:
        try:
            if hasattr(request.user, 'email') and request.user.email:
                user = UserProfile.objects.get(email=request.user.email)
                logger.info(f"✅ Found user by auth email: {user.email}")
                return user
        except UserProfile.DoesNotExist:
            pass
    
    # No fallback - user must provide user_id
    logger.warning(f"⚠️ No valid user_id provided and no authenticated user")
    return None


@api_view(['GET'])
@permission_classes([AllowAny])
def get_courses(request):
    """
    GET /api/trainee/courses/
    Get all available courses with enrollment info
    """
    try:
        from django.db import connection
        user = get_session_user(request)
        logger.info(f"Fetching courses for user: {user}")
        
        user_id = str(user.id) if user else None
        logger.info(f"📊 get_courses - User ID: {user_id}")
        
        # Use raw SQL to fetch courses with progress - ONLY enrolled courses
        cursor = connection.cursor()
        cursor.execute("""
            SELECT 
                c.course_id,
                c.title,
                c.description,
                c.course_type,
                c.is_mandatory,
                c.estimated_duration_hours,
                (SELECT COUNT(*) FROM modules WHERE course_id = c.course_id) as module_count,
                EXISTS(SELECT 1 FROM enrollments WHERE user_id = %s AND course_id = c.course_id) as enrolled,
                up.status,
                up.completion_percentage
            FROM courses c
            INNER JOIN enrollments e ON c.course_id = e.course_id AND e.user_id = %s
            LEFT JOIN user_progress up ON c.course_id = up.course_id AND up.user_id = %s
            ORDER BY c.title
        """, [user_id, user_id, user_id])
        
        course_data = []
        for row in cursor.fetchall():
            course_id, title, description, course_type, is_mandatory, duration, module_count, enrolled, course_status, completion_pct = row
            logger.info(f"📚 Course: {title}, Status: {course_status}, Progress: {completion_pct}%")
            course_data.append({
                'id': str(course_id),
                'course_id': str(course_id),
                'title': title,
                'description': description or '',
                'course_type': course_type,
                'is_mandatory': is_mandatory,
                'estimated_duration_hours': duration,
                'module_count': module_count or 0,
                'enrolled': enrolled,
                'status': course_status or 'not_started',
                'completion_percentage': completion_pct or 0,
            })
        
        cursor.close()
        logger.info(f"Found {len(course_data)} courses")
        
        return Response({
            'success': True,
            'courses': course_data,
            'count': len(course_data)
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error fetching courses: {str(e)}", exc_info=True)
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([AllowAny])
def get_course_detail(request, course_id):
    """
    GET /api/trainee/courses/{course_id}/
    Get detailed course information with modules
    """
    try:
        course = get_object_or_404(Course, course_id=course_id)
        
        # Get modules with basic info (avoid accessing related objects with schema issues)
        modules = course.modules.all().order_by('sequence_order')
        module_data = []
        
        for module in modules:
            module_info = {
                'id': str(module.module_id),
                'title': module.title,
                'description': module.description,
                'module_type': module.module_type,
                'position': module.sequence_order,
            }
            module_data.append(module_info)
        
        return Response({
            'success': True,
            'course': {
                'id': str(course.course_id),
                'title': course.title,
                'description': course.description,
                'course_type': course.course_type,
                'is_mandatory': course.is_mandatory,
                'estimated_duration_hours': course.estimated_duration_hours,
                'passing_criteria': course.passing_criteria,
                'modules': module_data,
                'module_count': len(module_data),
                'completion_percentage': 0,
            }
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error fetching course detail: {str(e)}")
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


def get_module_media(module):
    """Get all media associated with a module"""
    media_list = []
    
    # Videos - add directly as video_unit field for better frontend compatibility
    if hasattr(module, 'video_unit'):
        try:
            video = module.video_unit
            # Use video_url if available, otherwise construct path
            url = video.video_url or f'/media/{video.video_storage_path}' if video.video_storage_path else None
            if url:
                media_list.append({
                    'type': 'video',
                    'id': str(video.id),
                    'title': module.title,
                    'url': url,
                    'duration': video.duration,
                })
        except:
            pass
    
    # Audio
    if hasattr(module, 'audio_unit'):
        try:
            audio = module.audio_unit
            audio_path = audio.audio_storage_path.replace('\\', '/') if audio.audio_storage_path else None
            url = f'/media/{audio_path}' if audio_path else audio.audio_url
            media_list.append({
                'type': 'audio',
                'id': str(audio.id),
                'title': module.title,
                'url': url,
                'duration': audio.duration,
            })
        except:
            pass
    
    # Presentations
    if hasattr(module, 'presentation_unit'):
        try:
            presentation = module.presentation_unit
            pres_path = presentation.file_storage_path.replace('\\', '/') if presentation.file_storage_path else None
            url = f'/media/{pres_path}' if pres_path else presentation.file_url
            media_list.append({
                'type': 'presentation',
                'id': str(presentation.id),
                'title': module.title,
                'url': url,
            })
        except:
            pass
    
    # PDFs / Text
    if hasattr(module, 'text_unit'):
        try:
            text = module.text_unit
            text_path = text.file_storage_path.replace('\\', '/') if hasattr(text, 'file_storage_path') and text.file_storage_path else None
            url = f'/media/{text_path}' if text_path else (text.file_url if hasattr(text, 'file_url') else None)
            if url:
                media_list.append({
                    'type': 'pdf',
                    'id': str(text.id),
                    'title': module.title,
                    'url': url,
                })
        except:
            pass
    
    # SCORM
    if hasattr(module, 'scorm_package'):
        try:
            scorm = module.scorm_package
            media_list.append({
                'type': 'scorm',
                'id': str(scorm.id),
                'title': module.title,
                'url': f'/media/{scorm.file_storage_path}' if scorm.file_storage_path else scorm.file_url,
            })
        except:
            pass
    
    # Learning Resources (PDFs, PPTs, etc.)
    try:
        from trainee.models import LearningResource
        resources = module.resources.all().order_by('sequence_order')
        for resource in resources:
            # Map resource_type to frontend type
            resource_type_map = {
                'pdf': 'pdf',
                'ppt': 'ppt',
                'powerpoint': 'ppt',
                'docx': 'docx',
                'xlsx': 'xlsx',
                'image': 'image',
                'video': 'video',
                'link': 'link',
                'document': 'document',
            }
            frontend_type = resource_type_map.get(resource.resource_type, resource.resource_type)
            media_list.append({
                'type': frontend_type,
                'id': str(resource.resource_id),
                'title': resource.title,
                'url': resource.file_url,
                'description': resource.description,
                'file_size_bytes': resource.file_size_bytes,
            })
    except:
        pass
    
    return media_list


@api_view(['POST'])
@permission_classes([AllowAny])
def start_course(request, course_id):
    """
    POST /api/trainee/course/{course_id}/start/
    Start/enroll in a course
    """
    try:
        logger.info(f"Starting course: {course_id}")
        course = get_object_or_404(Course, course_id=course_id)
        logger.info(f"Course found: {course.title}")
        
        user = get_session_user(request)
        logger.info(f"User: {user}")
        
        if not user:
            logger.error("No user found for course start")
            return Response(
                {'success': False, 'error': 'No user found. Please log in.'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        logger.info(f"Creating enrollment for user {user} in course {course}")
        # Create or update enrollment
        enrollment, created = Enrollment.objects.get_or_create(
            user=user,
            course=course,
            defaults={'status': 'active'}
        )
        
        logger.info(f"Creating/updating progress")
        # Create or update progress
        progress, _ = UserProgress.objects.get_or_create(
            user=user,
            course=course,
            defaults={
                'status': 'in_progress',
                'completion_percentage': 0,
                'total_modules': course.modules.count(),
                'modules_completed': 0,
                'time_spent_minutes': 0,
            }
        )
        
        if progress.status == 'not_started':
            progress.status = 'in_progress'
            progress.save()
        
        logger.info(f"Course {course.title} started successfully for user {user}")
        return Response({
            'success': True,
            'message': 'Course started successfully',
            'course_id': str(course_id),
            'status': 'in_progress',
            'completion_percentage': progress.completion_percentage,
            'total_modules': progress.total_modules,
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error starting course: {str(e)}", exc_info=True)
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_400_BAD_REQUEST
        )


@api_view(['GET'])
@permission_classes([AllowAny])
def get_module_detail(request, module_id=None, course_id=None):
    """
    GET /api/trainee/modules/{module_id}/
    GET /api/trainee/course/{course_id}/modules/
    Get module details with media content or all modules for a course
    """
    try:
        # If course_id provided, get all modules for the course
        if course_id and not module_id:
            course = get_object_or_404(Course, course_id=course_id)
            modules = course.modules.all().order_by('created_at')
            
            module_data = []
            for module in modules:
                module_info = {
                    'id': str(module.module_id),
                    'title': module.title,
                    'description': module.description,
                    'module_type': module.module_type,
                    'position': module.sequence_order,
                    'media': get_module_media(module)
                }
                module_data.append(module_info)
            
            return Response({
                'success': True,
                'modules': module_data,
                'count': len(module_data)
            }, status=status.HTTP_200_OK)
        
        # Get specific module by ID
        module = get_object_or_404(Module, module_id=module_id)
        
        media = get_module_media(module)
        
        return Response({
            'success': True,
            'module': {
                'id': str(module.module_id),
                'title': module.title,
                'description': module.description,
                'module_type': module.module_type,
                'media': media,
            }
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error fetching module: {str(e)}")
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['GET'])
@permission_classes([AllowAny])
def get_video(request, video_id):
    """
    GET /api/trainee/videos/{video_id}/
    Get video details for streaming
    """
    try:
        video = get_object_or_404(VideoUnit, id=video_id)
        
        return Response({
            'success': True,
            'video': {
                'id': str(video.id),
                'title': video.unit.title,
                'description': video.unit.description,
                'url': f'/media/{video.video_storage_path}' if video.video_storage_path else video.video_url,
                'duration': video.duration,
                'module_id': str(video.unit.module_id),
                'course_id': str(video.unit.course.course_id),
            }
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error fetching video: {str(e)}")
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['POST'])
@permission_classes([AllowAny])
def mark_module_complete(request, module_id):
    """
    POST /api/trainee/modules/{module_id}/complete/
    Mark module as completed
    """
    try:
        module = get_object_or_404(Module, module_id=module_id)
        user = get_session_user(request)
        
        if not user:
            return Response(
                {'success': False, 'error': 'No user found'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Create completion record
        completion, created = ModuleCompletion.objects.get_or_create(
            user=user,
            module=module,
        )
        
        # Update course progress
        course = module.course
        total_modules = course.modules.count()
        completed_modules = ModuleCompletion.objects.filter(
            user=user,
            module__course=course
        ).count()
        
        completion_percentage = (completed_modules / total_modules * 100) if total_modules > 0 else 0
        
        progress, _ = UserProgress.objects.get_or_create(
            user=user,
            course=course,
            defaults={'status': 'in_progress'}
        )
        
        progress.modules_completed = completed_modules
        progress.completion_percentage = completion_percentage
        if completion_percentage == 100:
            progress.status = 'completed'
        progress.save()
        
        return Response({
            'success': True,
            'message': 'Module marked as complete',
            'module_id': str(module_id),
            'course_progress': {
                'completion_percentage': completion_percentage,
                'modules_completed': completed_modules,
                'total_modules': total_modules,
            }
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error marking module complete: {str(e)}")
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_400_BAD_REQUEST
        )


@api_view(['GET'])
@permission_classes([AllowAny])
def get_dashboard(request):
    """
    GET /api/trainee/dashboard/
    Get user dashboard with all courses and progress
    """
    try:
        from django.db import connection
        user = get_session_user(request)
        
        if not user:
            return Response({
                'success': False,
                'error': 'No user found'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        user_id = str(user.id)  # UserProfile.id maps to user_id column
        logger.info(f"📊 get_dashboard - User ID: {user_id}")
        
        # Get all courses with progress using raw SQL - ONLY enrolled courses
        cursor = connection.cursor()
        cursor.execute("""
            SELECT 
                c.course_id,
                c.title,
                c.description,
                up.status,
                up.completion_percentage,
                (SELECT COUNT(*) FROM modules WHERE course_id = c.course_id) as module_count
            FROM courses c
            INNER JOIN enrollments e ON c.course_id = e.course_id AND e.user_id = %s
            LEFT JOIN user_progress up ON c.course_id = up.course_id AND up.user_id = %s
            ORDER BY c.title
        """, [user_id, user_id])
        
        courses_data = {
            'total': 0,
            'in_progress': 0,
            'not_started': 0,
            'completed': 0,
            'courses': []
        }
        
        for row in cursor.fetchall():
            course_id, title, description, status_val, completion_pct, module_count = row
            
            # Default to not_started if no progress record
            status_val = status_val or 'not_started'
            completion_pct = completion_pct or 0
            
            logger.info(f"📚 Dashboard - Course: {title}, Status: {status_val}, Progress: {completion_pct}%")
            
            courses_data['total'] += 1
            if status_val == 'in_progress':
                courses_data['in_progress'] += 1
            elif status_val == 'completed':
                courses_data['completed'] += 1
            else:
                courses_data['not_started'] += 1
            
            courses_data['courses'].append({
                'id': str(course_id),
                'title': title,
                'description': description or '',
                'status': status_val,
                'completion_percentage': completion_pct,
                'module_count': module_count or 0,
            })
        
        cursor.close()
        
        return Response({
            'success': True,
            'user': {
                'name': f'{user.first_name} {user.last_name}',
                'email': user.email,
                'role': user.role,
            },
            'dashboard': courses_data
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error fetching dashboard: {str(e)}", exc_info=True)
        import traceback
        traceback.print_exc()
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import csrf_exempt

@csrf_exempt
@require_http_methods(["GET"])
def convert_ppt_to_pdf(request, resource_id):
    """
    GET /api/trainee/convert-ppt/{resource_id}/
    Convert PPT file to PDF using python-pptx library
    Streams binary PDF response (not JSON)
    """
    try:
        from trainee.models import LearningResource, MediaMetadata
        from trainee.services.ppt_converter import PPTToPDFConverter
        import os
        import hashlib
        from django.conf import settings
        
        logger.info(f"=== PPT to PDF Conversion Request ===")
        logger.info(f"Resource ID: {resource_id}")
        
        from django.http import JsonResponse, HttpResponse
        
        # Get the resource
        try:
            resource = LearningResource.objects.get(resource_id=resource_id)
            logger.info(f"✓ Found LearningResource: {resource.title}")
            logger.info(f"  Type: {resource.resource_type}")
            logger.info(f"  File URL: {resource.file_url}")
        except LearningResource.DoesNotExist:
            logger.error(f"✗ LearningResource not found: {resource_id}")
            return JsonResponse(
                {'success': False, 'error': f'Resource not found: {resource_id}'},
                status=404
            )
        
        # Check if it's a PPT file
        if resource.resource_type not in ['ppt', 'powerpoint', 'presentation']:
            return JsonResponse(
                {'success': False, 'error': f'Resource is not a PowerPoint file (type: {resource.resource_type})'},
                status=400
            )
        
        # Get file path from resource.file_url
        file_url = resource.file_url
        logger.info(f"Processing file_url: {file_url}")
        
        if file_url.startswith('/media/'):
            file_path = os.path.join(settings.MEDIA_ROOT, file_url[7:])
        else:
            file_path = os.path.join(settings.MEDIA_ROOT, file_url)
        
        # Normalize path
        file_path = os.path.normpath(file_path)
        logger.info(f"Computed file path: {file_path}")
        logger.info(f"MEDIA_ROOT: {settings.MEDIA_ROOT}")
        logger.info(f"File exists: {os.path.exists(file_path)}")
        
        # If file doesn't exist, try to find it from MediaMetadata using storage_path
        if not os.path.exists(file_path):
            logger.warning(f"File not found at {file_path}, attempting to locate via MediaMetadata...")
            
            # Try to find the file in MediaMetadata table by matching filename
            filename = os.path.basename(file_path)
            logger.info(f"Searching for file: {filename}")
            
            try:
                # Search for file by name or storage_path pattern
                media_files = MediaMetadata.objects.filter(
                    file_name__icontains=os.path.splitext(filename)[0]
                )
                logger.info(f"Found {media_files.count()} matching media files")
                
                if media_files.exists():
                    media_file = media_files.first()
                    storage_path = media_file.storage_path
                    logger.info(f"✓ Found in MediaMetadata: {storage_path}")
                    
                    # Construct file path from storage_path
                    file_path = os.path.join(settings.MEDIA_ROOT, storage_path)
                    file_path = os.path.normpath(file_path)
                    logger.info(f"Updated file path: {file_path}")
                    logger.info(f"File exists now: {os.path.exists(file_path)}")
                else:
                    logger.warning(f"No matching MediaMetadata found for {filename}")
            except Exception as e:
                logger.warning(f"Error searching MediaMetadata: {str(e)}")
        
        # Final check
        if not os.path.exists(file_path):
            logger.error(f"✗ PPT file not found: {file_path}")
            return JsonResponse({
                'success': False, 
                'error': f'File not found on server',
                'debug': {
                    'file_url': resource.file_url,
                    'computed_path': file_path,
                    'media_root': settings.MEDIA_ROOT
                }
            }, status=404)
        
        logger.info(f"✓ PPT file found and ready for conversion")
        
        # Convert PPT to PDF in-memory and stream it back (no disk PDF creation)
        logger.info("Starting in-memory PPT to PDF conversion...")
        try:
            import io
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from pptx import Presentation
            from django.http import HttpResponse

            prs = Presentation(file_path)
            buffer = io.BytesIO()
            pdf_canvas = canvas.Canvas(buffer, pagesize=letter)
            width, height = letter

            slide_count = len(prs.slides)
            logger.info(f"Converting {slide_count} slides in-memory")

            for slide_idx, slide in enumerate(prs.slides, 1):
                try:
                    # Extract text from slide
                    text_content = []
                    try:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text') and shape.text and shape.text.strip():
                                text_content.append(shape.text.strip())
                    except Exception:
                        logger.debug(f"Could not extract shape text for slide {slide_idx}")

                    slide_text = "\n\n".join(text_content)

                    # Add page header
                    pdf_canvas.showPage()
                    pdf_canvas.setFont('Helvetica-Bold', 14)
                    pdf_canvas.drawString(40, height - 40, f"Slide {slide_idx}")
                    pdf_canvas.setStrokeColorRGB(0.7, 0.7, 0.7)
                    pdf_canvas.line(40, height - 50, width - 40, height - 50)

                    if slide_text:
                        pdf_canvas.setFont('Helvetica', 10)
                        y_position = height - 80
                        lines = slide_text.split('\n')
                        for line in lines:
                            if y_position < 40:
                                pdf_canvas.showPage()
                                pdf_canvas.setFont('Helvetica', 10)
                                y_position = height - 40

                            # Simple wrap for very long lines
                            if len(line) > 100:
                                chunks = [line[i:i+100] for i in range(0, len(line), 100)]
                                for chunk in chunks:
                                    pdf_canvas.drawString(50, y_position, chunk)
                                    y_position -= 15
                            else:
                                pdf_canvas.drawString(50, y_position, line)
                                y_position -= 15
                    else:
                        pdf_canvas.setFont('Helvetica', 12)
                        pdf_canvas.setFillColorRGB(0.8, 0.8, 0.8)
                        pdf_canvas.drawString(40, height / 2, "[Slide content]")

                except Exception as e:
                    logger.warning(f"Error processing slide {slide_idx}: {str(e)}")
                    pdf_canvas.showPage()
                    pdf_canvas.setFont('Helvetica-Bold', 14)
                    pdf_canvas.drawString(40, height - 40, f"Slide {slide_idx} (error)")

            pdf_canvas.save()
            buffer.seek(0)

            filename = os.path.splitext(os.path.basename(file_path))[0] + '.pdf'
            response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
            response['Content-Disposition'] = f'inline; filename="{filename}"'
            logger.info(f"Streaming PDF for resource {resource_id} as {filename}")
            return response

        except Exception as e:
            logger.error(f"✗ In-memory PPT conversion failed: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return JsonResponse({
                'success': False,
                'error': f'In-memory conversion failed: {str(e)}'
            }, status=500)
    
    except Exception as e:
        logger.error(f"✗ Unexpected error in convert_ppt_to_pdf: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        return JsonResponse(
            {'success': False, 'error': f'Server error: {str(e)}'},
            status=500
        )



@api_view(['GET'])
@permission_classes([AllowAny])
def get_module_mixed_content(request, module_id):
    """
    GET /api/trainee/module/{module_id}/content/
    Get all content for a module (videos, PDFs, PPTs, quizzes, etc.)
    """
    try:
        module = get_object_or_404(Module, module_id=module_id)
        user = get_session_user(request)
        
        content_list = []
        
        # First, check if media_metadata has content for this module
        has_media_metadata = False
        try:
            from django.db import connection
            cursor = connection.cursor()
            cursor.execute('SELECT COUNT(*) FROM media_metadata WHERE unit_id = %s', [str(module.module_id)])
            media_count = cursor.fetchone()[0]
            has_media_metadata = media_count > 0
        except Exception as e:
            logger.error(f"Error checking media_metadata: {str(e)}")
        
        # If media_metadata has content, use ONLY that (to avoid duplicates)
        if has_media_metadata:
            try:
                from django.db import connection
                cursor = connection.cursor()
                cursor.execute('''
                    SELECT media_id, file_name, file_type, mime_type, storage_path, file_size, duration
                    FROM media_metadata 
                    WHERE unit_id = %s
                ''', [str(module.module_id)])
                
                for row in cursor.fetchall():
                    media_id, file_name, file_type, mime_type, storage_path, file_size, duration = row
                    
                    # Build file URL with full domain
                    if storage_path:
                        file_url = request.build_absolute_uri(f'/media/{storage_path}')
                        
                        # Determine content type from file_type
                        if file_type == 'video' or (mime_type and 'video' in mime_type):
                            content_type = 'video'
                        elif file_type == 'pdfs' or (mime_type and 'pdf' in mime_type):
                            content_type = 'pdf'
                        elif file_type in ['ppt', 'pptx'] or (mime_type and 'presentation' in mime_type):
                            content_type = 'ppt'
                        else:
                            content_type = 'document'
                        
                        content_item = {
                            'id': str(media_id),
                            'media_id': str(media_id),
                            'content_type': content_type,
                            'title': module.title,
                            'description': f'{content_type.upper()} for {module.title}',
                            'file_url': file_url,
                            'file_name': file_name,
                            'is_unlocked': True,
                        }
                        
                        # Add type-specific fields
                        if content_type == 'video' and duration:
                            content_item['duration'] = duration
                            content_item['duration_seconds'] = duration
                        if file_size:
                            content_item['file_size_mb'] = file_size / (1024 * 1024)
                        
                        content_list.append(content_item)
                        logger.info(f"Added media from media_metadata: {file_name} ({content_type})")
            except Exception as e:
                logger.error(f"Error adding media from media_metadata: {str(e)}")
        else:
            # Fallback: Use video_units and presentation_units tables if no media_metadata
            # 1. Add Video Unit
            try:
                if hasattr(module, 'video_unit') and module.video_unit:
                    video = module.video_unit
                    # Construct video URL with full domain: prefer video_url, fall back to video_storage_path
                    video_file_url = video.video_url
                    if not video_file_url and video.video_storage_path:
                        # Normalize path - convert backslashes to forward slashes
                        normalized_path = video.video_storage_path.replace('\\', '/')
                        video_file_url = request.build_absolute_uri(f'/media/{normalized_path}')
                    
                    if video_file_url:
                        content_list.append({
                            'id': str(video.id),
                            'content_type': 'video',
                            'title': module.title,
                            'description': f'Video lesson for {module.title}',
                            'file_url': video_file_url,
                            'duration': video.duration,
                            'duration_seconds': video.duration,
                            'is_unlocked': True,
                        })
                        logger.info(f"Added video for module {module.module_id}: {video_file_url}")
                    else:
                        logger.warning(f"No video URL or storage path for module {module.module_id}")
            except Exception as e:
                logger.error(f"Error adding video unit for module {module.module_id}: {str(e)}")
                pass
            
            # 2. Add Presentation Unit (PPT/PDF from PostgreSQL presentation_units table)
            try:
                if hasattr(module, 'presentation_unit') and module.presentation_unit:
                    presentation = module.presentation_unit
                    # Build presentation file URL with full domain
                    ppt_url = presentation.file_url
                    if not ppt_url and presentation.file_storage_path:
                        # Normalize path
                        normalized_path = presentation.file_storage_path.replace('\\', '/')
                        ppt_url = request.build_absolute_uri(f'/media/{normalized_path}')
                    
                    if ppt_url:
                        # Determine file type from extension
                        file_extension = ppt_url.lower().split('.')[-1] if '.' in ppt_url else 'pdf'
                        content_type = 'pdf' if file_extension == 'pdf' else 'ppt'
                        
                        content_list.append({
                            'id': str(presentation.id),
                            'presentation_id': str(presentation.id),
                            'content_type': content_type,
                            'title': module.title,
                            'description': f'Presentation for {module.title}',
                            'file_url': ppt_url,
                            'slide_count': presentation.slide_count if hasattr(presentation, 'slide_count') else None,
                            'is_unlocked': True,
                        })
                        logger.info(f"Added presentation for module {module.module_id}: {ppt_url}")
                    else:
                        logger.warning(f"No presentation URL or storage path for module {module.module_id}")
            except Exception as e:
                logger.error(f"Error adding presentation unit for module {module.module_id}: {str(e)}")
                pass
        
        # 3. Add Learning Resources (PDFs, PPTs, etc.)
        try:
            from trainee.models import LearningResource
            resources = module.resources.all().order_by('sequence_order')
            for resource in resources:
                resource_type = resource.resource_type.lower() if resource.resource_type else 'resource'
                content_list.append({
                    'id': str(resource.resource_id),
                    'resource_id': str(resource.resource_id),
                    'content_type': 'resource',
                    'resource_type': resource_type,
                    'title': resource.title,
                    'description': resource.description or f'{resource_type.upper()} resource',
                    'file_url': resource.file_url,
                    'file_size_mb': resource.file_size_bytes / (1024 * 1024) if resource.file_size_bytes else None,
                    'is_unlocked': True,
                })
        except Exception as e:
            logger.error(f"Error adding learning resources: {str(e)}")
            pass
        
        # 3. Add Quizzes (from new Quiz table)
        try:
            from trainee.models import Quiz
            quizzes = Quiz.objects.filter(unit=module)
            for quiz in quizzes:
                questions_count = quiz.questions.count() if hasattr(quiz, 'questions') else 0
                content_list.append({
                    'id': str(quiz.id),
                    'quiz_id': str(quiz.id),
                    'content_type': 'quiz',
                    'title': f'Quiz: {module.title}',
                    'description': f'Quiz for {module.title}',
                    'time_limit_minutes': quiz.time_limit,
                    'passing_score': quiz.passing_score,
                    'points_possible': 100,
                    'questions_count': questions_count,
                    'is_unlocked': True,
                    'attempts_allowed': quiz.attempts_allowed,
                    'show_answers': quiz.show_answers,
                })
        except Exception as e:
            logger.error(f"Error adding quizzes: {str(e)}")
            pass
        
        return Response({
            'success': True,
            'content': content_list,
            'module_id': str(module.module_id),
            'module_title': module.title,
        }, status=status.HTTP_200_OK)
        
    except Exception as e:
        logger.error(f"Error fetching module mixed content: {str(e)}")
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )
    except Exception as e:
        logger.error(f"Error in convert_ppt_to_pdf: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        return Response(
            {'success': False, 'error': str(e)},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )

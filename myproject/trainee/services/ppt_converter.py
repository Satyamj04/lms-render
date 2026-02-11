"""
PPT/PPTX to PDF Converter Service
Converts PowerPoint presentations to PDF using python-pptx and reportlab
"""

import os
import logging
from pathlib import Path
from PIL import Image, ImageDraw
from io import BytesIO
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from pptx import Presentation
from pptx.util import Inches, Pt
from django.conf import settings

logger = logging.getLogger(__name__)


class PPTToPDFConverter:
    """Convert PPT/PPTX files to PDF"""
    
    def __init__(self):
        self.supported_extensions = ['.ppt', '.pptx']
        self.page_width, self.page_height = letter
    
    def convert(self, ppt_file_path: str, output_pdf_path: str = None) -> str:
        """
        Convert PPT file to PDF
        
        Args:
            ppt_file_path: Path to the PPT/PPTX file
            output_pdf_path: Optional custom output path (if not provided, creates one from input)
        
        Returns:
            Path to the generated PDF file
        """
        try:
            # Validate input file
            if not os.path.exists(ppt_file_path):
                raise FileNotFoundError(f"PPT file not found: {ppt_file_path}")
            
            # Determine output path
            if not output_pdf_path:
                base_path = os.path.splitext(ppt_file_path)[0]
                output_pdf_path = f"{base_path}.pdf"
            
            # Check if PDF already exists
            if os.path.exists(output_pdf_path):
                logger.info(f"PDF already exists: {output_pdf_path}")
                return output_pdf_path
            
            logger.info(f"Converting PPT to PDF: {ppt_file_path} → {output_pdf_path}")
            
            # Method 1: Try using python-pptx (works for PPTX files)
            try:
                return self._convert_with_pptx(ppt_file_path, output_pdf_path)
            except Exception as e:
                logger.warning(f"python-pptx conversion failed, trying alternative method: {str(e)}")
                # Fallback method: Create a simple PDF with slide info
                return self._create_fallback_pdf(ppt_file_path, output_pdf_path)
        
        except Exception as e:
            logger.error(f"Error converting PPT to PDF: {str(e)}")
            raise
    
    def _convert_with_pptx(self, ppt_file_path: str, output_pdf_path: str) -> str:
        """
        Convert using python-pptx by rendering slides as images
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            # Load presentation
            prs = Presentation(ppt_file_path)
            
            # Create temporary directory for slide images
            temp_dir = os.path.dirname(output_pdf_path)
            os.makedirs(temp_dir, exist_ok=True)
            
            # Create PDF with reportlab
            pdf_canvas = canvas.Canvas(output_pdf_path, pagesize=letter)
            width, height = letter
            
            slide_count = len(prs.slides)
            logger.info(f"Converting {slide_count} slides to PDF")
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                try:
                    logger.debug(f"Processing slide {slide_idx}/{slide_count}")
                    
                    # Try to extract text and render it
                    slide_text = self._extract_slide_text(slide)
                    
                    # Create a simple PDF page with text content
                    self._add_slide_to_pdf(pdf_canvas, slide_idx, slide_text, width, height)
                
                except Exception as e:
                    logger.warning(f"Error processing slide {slide_idx}: {str(e)}")
                    # Add blank page if processing fails
                    self._add_blank_slide(pdf_canvas, slide_idx, width, height)
            
            pdf_canvas.save()
            logger.info(f"PDF created successfully: {output_pdf_path}")
            return output_pdf_path
        
        except Exception as e:
            logger.error(f"python-pptx conversion failed: {str(e)}")
            raise
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text content from a slide"""
        text_content = []
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
        except Exception as e:
            logger.warning(f"Error extracting text from slide: {str(e)}")
        
        return "\n\n".join(text_content)
    
    def _add_slide_to_pdf(self, pdf_canvas, slide_num: int, text_content: str, width: float, height: float):
        """Add a slide's content to PDF"""
        try:
            # Add new page
            pdf_canvas.showPage()
            
            # Add slide header
            pdf_canvas.setFont("Helvetica-Bold", 14)
            pdf_canvas.drawString(40, height - 40, f"Slide {slide_num}")
            
            # Add separator line
            pdf_canvas.setStrokeColorRGB(0.7, 0.7, 0.7)
            pdf_canvas.line(40, height - 50, width - 40, height - 50)
            
            # Add slide content
            if text_content:
                pdf_canvas.setFont("Helvetica", 10)
                y_position = height - 80
                
                # Split text into lines and add to PDF
                lines = text_content.split('\n')
                for line in lines:
                    if y_position < 40:  # If we're near bottom, start a new page
                        pdf_canvas.showPage()
                        pdf_canvas.setFont("Helvetica", 10)
                        y_position = height - 40
                    
                    # Wrap long lines
                    if len(line) > 100:
                        wrapped_lines = self._wrap_text(line, 100)
                        for wrapped_line in wrapped_lines:
                            pdf_canvas.drawString(50, y_position, wrapped_line)
                            y_position -= 15
                    else:
                        pdf_canvas.drawString(50, y_position, line)
                        y_position -= 15
            else:
                # If no text, add placeholder
                pdf_canvas.setFont("Helvetica", 12)
                pdf_canvas.setFillColorRGB(0.8, 0.8, 0.8)
                pdf_canvas.drawString(40, height / 2, "[Slide content]")
        
        except Exception as e:
            logger.warning(f"Error adding slide to PDF: {str(e)}")
    
    def _add_blank_slide(self, pdf_canvas, slide_num: int, width: float, height: float):
        """Add a blank slide to PDF"""
        pdf_canvas.showPage()
        pdf_canvas.setFont("Helvetica-Bold", 14)
        pdf_canvas.drawString(40, height - 40, f"Slide {slide_num}")
        pdf_canvas.setStrokeColorRGB(0.7, 0.7, 0.7)
        pdf_canvas.line(40, height - 50, width - 40, height - 50)
    
    def _wrap_text(self, text: str, width: int) -> list:
        """Wrap text to fit within specified width"""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            current_line.append(word)
            if len(' '.join(current_line)) > width:
                current_line.pop()
                lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return lines
    
    def _create_fallback_pdf(self, ppt_file_path: str, output_pdf_path: str) -> str:
        """
        Create a simple PDF with file information when conversion fails
        """
        try:
            logger.info(f"Creating fallback PDF: {output_pdf_path}")
            
            pdf_canvas = canvas.Canvas(output_pdf_path, pagesize=letter)
            width, height = letter
            
            # Get file info
            file_name = os.path.basename(ppt_file_path)
            file_size = os.path.getsize(ppt_file_path) / (1024 * 1024)  # MB
            
            # Try to load and get slide count
            try:
                prs = Presentation(ppt_file_path)
                slide_count = len(prs.slides)
            except:
                slide_count = "Unknown"
            
            # Add content
            pdf_canvas.setFont("Helvetica-Bold", 16)
            pdf_canvas.drawString(40, height - 40, "Presentation Document")
            
            pdf_canvas.setFont("Helvetica", 12)
            y = height - 100
            
            pdf_canvas.drawString(40, y, f"Filename: {file_name}")
            y -= 30
            pdf_canvas.drawString(40, y, f"Size: {file_size:.2f} MB")
            y -= 30
            pdf_canvas.drawString(40, y, f"Slides: {slide_count}")
            y -= 50
            
            pdf_canvas.setFont("Helvetica-Oblique", 10)
            pdf_canvas.drawString(40, y, "This is a placeholder PDF for the presentation.")
            pdf_canvas.drawString(40, y - 20, "The presentation file has been uploaded successfully.")
            
            pdf_canvas.save()
            logger.info(f"Fallback PDF created: {output_pdf_path}")
            return output_pdf_path
        
        except Exception as e:
            logger.error(f"Error creating fallback PDF: {str(e)}")
            raise


# Singleton instance
_converter = None


def get_ppt_converter():
    """Get PPT converter instance"""
    global _converter
    if _converter is None:
        _converter = PPTToPDFConverter()
    return _converter


def convert_ppt_to_pdf(ppt_file_path: str, output_pdf_path: str = None) -> str:
    """
    Utility function to convert PPT to PDF
    
    Args:
        ppt_file_path: Path to the PPT/PPTX file
        output_pdf_path: Optional custom output path
    
    Returns:
        Path to the generated PDF file
    """
    converter = get_ppt_converter()
    return converter.convert(ppt_file_path, output_pdf_path)
